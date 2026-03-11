#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import re
import tempfile
import time
import urllib.request
from copy import deepcopy
from datetime import datetime

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patheffects as path_effects
from matplotlib.ticker import MaxNLocator, ScalarFormatter
import numpy as np
import pandas as pd
from matplotlib import font_manager
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn


TEMPLATE = '1.26，模版 - 07022025.pptx'
EXCEL = '工作簿5.xlsx'
OUTPUT = '生成结果.pptx'
CHART_PATH = '市值_PS_TTM_可视化.png'


def find_columns(raw):
    issue_col = None
    spac_col = None
    for i, v in enumerate(raw.iloc[1]):
        if isinstance(v, str) and '发行方式' in v:
            issue_col = i
        if isinstance(v, str) and '是否SPAC上市' in v:
            spac_col = i

    logo_col = None
    url_re = re.compile(r'^https?://', re.I)
    for i in range(raw.shape[1]):
        v = raw.iloc[1, i]
        if isinstance(v, str) and url_re.match(v) and 'logo' in v.lower():
            logo_col = i
            break
    if logo_col is None:
        for i in range(raw.shape[1]):
            col_vals = raw.iloc[2:, i].dropna()
            if col_vals.empty:
                continue
            sample = col_vals.iloc[0]
            if isinstance(sample, str) and url_re.match(sample):
                logo_col = i
                break
    if logo_col is None:
        for i, v in enumerate(raw.iloc[0]):
            if isinstance(v, str) and 'logo' in v.lower():
                logo_col = i
                break

    return issue_col, spac_col, logo_col


def fmt(v):
    if pd.isna(v):
        return ''
    if isinstance(v, (pd.Timestamp, datetime)):
        return v.strftime('%Y-%m-%d')
    if isinstance(v, (int, float, np.number)) and not isinstance(v, bool):
        return f"{float(v):.2f}"
    return str(v)


def load_rows(excel_path):
    raw = pd.read_excel(excel_path, header=None)
    keys = list(raw.iloc[0])
    col_keys = {i: k for i, k in enumerate(keys) if isinstance(k, str) and k.strip()}

    issue_col, spac_col, logo_col = find_columns(raw)

    rows = []
    data_rows = raw.iloc[2:].reset_index(drop=True)
    for _, row in data_rows.iterrows():
        if row.isna().all():
            continue
        row_dict = {}
        for col, key in col_keys.items():
            row_dict[key] = fmt(row[col])
        row_dict['issue_method'] = fmt(row[issue_col]) if issue_col is not None else ''
        row_dict['spac_flag'] = fmt(row[spac_col]) if spac_col is not None else ''
        row_dict['logo_url'] = fmt(row[logo_col]) if logo_col is not None else ''

        mcap_raw = row_dict.get('market_cap', '')
        try:
            mcap_val = float(mcap_raw)
            row_dict['market_cap_level'] = '大' if mcap_val > 100 else '小'
        except Exception:
            row_dict['market_cap_level'] = ''

        rows.append(row_dict)

    return rows


def pick_cjk_font():
    preferred_fonts = ['STHeiti', 'SimHei', 'PingFang SC', 'Heiti SC', 'Songti SC', 'Microsoft YaHei']
    for f in font_manager.fontManager.ttflist:
        if f.name in preferred_fonts:
            return f.name
    return None


def build_chart(rows, chart_path):
    labels, mcaps, ps_vals = [], [], []
    for r in rows:
        try:
            m = float(r.get('market_cap', ''))
            p = float(r.get('ps_ttm', ''))
        except Exception:
            continue
        label = r.get('company_name') or r.get('ticker') or ''
        labels.append(label)
        mcaps.append(m)
        ps_vals.append(p)

    if not labels:
        return

    # Use a unified Chinese font (SimHei) if available; fallback to any CJK font.
    font_name = 'STHeiti'
    available = {f.name for f in font_manager.fontManager.ttflist}
    if font_name not in available:
        font_name = pick_cjk_font()
    if font_name:
        plt.rcParams['font.sans-serif'] = [font_name]
        plt.rcParams['font.family'] = 'sans-serif'
        plt.rcParams['axes.unicode_minus'] = False
        font_props = {'fontfamily': font_name}
    else:
        font_props = {}

    color_bar = '#7A0019'
    color_line = '#000000'
    color_avg = '#000000'
    bold_effects = [path_effects.Stroke(linewidth=1.2, foreground='#000000'), path_effects.Normal()]
    bar_effects = [path_effects.Stroke(linewidth=1.2, foreground=color_bar), path_effects.Normal()]

    x = np.arange(len(labels))
    fig, ax1 = plt.subplots(figsize=(12, 7), dpi=120)

    bars = ax1.bar(x, mcaps, color=color_bar, alpha=0.85, width=0.5, label='总市值 (亿)')
    ax1.set_ylabel('总市值 (亿元)', color=color_bar, fontsize=13, fontweight='bold')
    ax1.tick_params(axis='y', labelcolor=color_bar)
    ax1.set_xticks(x)
    ax1.set_xticklabels(labels)
    ax1.set_ylim(0, max(mcaps) * 1.15)
    ax1.yaxis.set_major_locator(MaxNLocator(nbins=5))
    ax1.yaxis.set_major_formatter(ScalarFormatter(useMathText=False))
    ax1.ticklabel_format(style='plain', axis='y')

    for rect, val in zip(bars, mcaps):
        ax1.text(
            rect.get_x() + rect.get_width() / 2,
            rect.get_height() + (max(mcaps) * 0.03),
            f'{val:.2f}',
            ha='center',
            va='bottom',
            fontsize=11,
            color=color_bar,
            fontweight='bold',
            **font_props,
            path_effects=bar_effects,
        )

    ax2 = ax1.twinx()
    ax2.scatter(x, ps_vals, color=color_line, marker='o', s=35, label='市销率 PS (TTM)')
    ps_avg = sum(ps_vals) / len(ps_vals)
    ax2.axhline(y=ps_avg, color=color_avg, linestyle='--', linewidth=2, alpha=0.8)
    ax2.set_ylabel('市销率 PS (倍)', color=color_line, fontsize=13, fontweight='bold')
    ax2.tick_params(axis='y', labelcolor=color_line)
    ax2.set_ylim(0, max(ps_vals) * 1.2)
    ax2.yaxis.set_major_locator(MaxNLocator(nbins=5))
    ax2.yaxis.set_major_formatter(ScalarFormatter(useMathText=False))
    ax2.ticklabel_format(style='plain', axis='y')

    for i, val in enumerate(ps_vals):
        # Stagger PS labels and shift slightly in x to avoid overlap with bar labels.
        dx = 0.18 if i % 2 == 0 else -0.18
        dy = (0.05 + (i % 2) * 0.03) * max(ps_vals)
        ax2.text(
            i + dx,
            val + dy,
            f'{val:.2f}',
            ha='center',
            va='bottom',
            fontsize=11,
            color='#000000',
            fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.15', facecolor='white', edgecolor='none', alpha=0.6),
            **font_props,
            path_effects=bold_effects,
        )

    plt.title('各公司市值与PS估值对比', fontsize=17, fontweight='bold', pad=20, color='#333333')

    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax2.legend(
        lines1 + lines2,
        labels1 + labels2,
        loc='upper center',
        bbox_to_anchor=(0.5, -0.1),
        ncol=2,
        frameon=False,
        prop={'size': 11, 'weight': 'bold'},
    )

    # Ensure tick labels are bold.
    for lbl in ax1.get_xticklabels() + ax1.get_yticklabels() + ax2.get_yticklabels():
        lbl.set_fontweight('bold')

    plt.tight_layout()
    plt.savefig(chart_path, bbox_inches='tight')
    plt.close()


def clone_slide(prs, src_slide):
    layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for shape in list(new_slide.shapes):
        if shape.is_placeholder:
            el = shape.element
            el.getparent().remove(el)
    rel_map = {}
    for rel in src_slide.part.rels.values():
        if rel.reltype == RT.IMAGE:
            new_rId = new_slide.part.rels._add_relationship(rel.reltype, rel._target, is_external=False)
            rel_map[rel.rId] = new_rId
    for shape in src_slide.shapes:
        new_el = deepcopy(shape.element)
        for blip in new_el.xpath('.//a:blip'):
            rId = blip.get(qn('r:embed'))
            if rId in rel_map:
                blip.set(qn('r:embed'), rel_map[rId])
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide


def replace_text_in_slide(slide, mapping):
    for shape in list(slide.shapes):
        if not hasattr(shape, 'text_frame'):
            continue
        tf = shape.text_frame
        for paragraph in tf.paragraphs:
            texts = paragraph._p.xpath('.//a:t')
            if not texts:
                continue
            full = ''.join(t.text or '' for t in texts)
            new = full
            for k, v in mapping.items():
                new = new.replace('{{' + k + '}}', v)
            if new == full:
                continue
            remaining = new
            for i, t in enumerate(texts):
                if i == len(texts) - 1:
                    t.text = remaining
                else:
                    orig = t.text or ''
                    take = len(orig)
                    t.text = remaining[:take]
                    remaining = remaining[take:]
        if not tf.text.strip():
            el = shape.element
            el.getparent().remove(el)


def replace_logo(slide, logo_url, cache):
    if not logo_url:
        return
    def download_logo(url, dest, attempts=3, timeout=10):
        last_exc = None
        for i in range(attempts):
            try:
                with urllib.request.urlopen(url, timeout=timeout) as resp:
                    data = resp.read()
                with open(dest, 'wb') as f:
                    f.write(data)
                return True
            except Exception as e:
                last_exc = e
                print(f'logo download failed ({i+1}/{attempts}): {url} - {e}')
                time.sleep(1.0 * (i + 1))
        return False
    try:
        if logo_url in cache:
            img_path = cache[logo_url]
        else:
            fd, img_path = tempfile.mkstemp(suffix='.png')
            os.close(fd)
            ok = download_logo(logo_url, img_path)
            if not ok:
                try:
                    os.remove(img_path)
                except Exception:
                    pass
                return
            cache[logo_url] = img_path
    except Exception:
        return

    pic_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pic_shape = shape
            break
    if not pic_shape:
        return

    left, top, width, height = pic_shape.left, pic_shape.top, pic_shape.width, pic_shape.height
    el = pic_shape.element
    el.getparent().remove(el)
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)


def replace_summary_image(slide, chart_path):
    if not os.path.exists(chart_path):
        return
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            el = shape.element
            el.getparent().remove(el)
            slide.shapes.add_picture(chart_path, left, top, width=width, height=height)
            break


def generate_ppt(
    excel_data_path=None,
    output_path=None,
    template_path=None,
    chart_path=None,
    rows=None,
):
    """Generate PPTX from excel file or pre-parsed rows."""
    if rows is None:
        if not excel_data_path:
            raise ValueError('excel_data_path is required when rows is not provided')
        rows = load_rows(excel_data_path)

    if not template_path:
        template_path = TEMPLATE
    if not output_path:
        output_path = OUTPUT

    cleanup_chart = False
    if chart_path is None:
        fd, chart_path = tempfile.mkstemp(suffix='.png')
        os.close(fd)
        cleanup_chart = True

    build_chart(rows, chart_path)

    prs = Presentation(template_path)
    original_slide_count = len(prs.slides)
    slide_cover = prs.slides[0]
    slide_ipo = prs.slides[1]
    slide_spac = prs.slides[2]
    slide_summary = prs.slides[3]
    slide_thanks = prs.slides[4]

    cache = {}

    clone_slide(prs, slide_cover)

    for row in rows:
        spac_flag = row.get('spac_flag', '')
        use_spac = False
        if isinstance(spac_flag, str):
            val = spac_flag.strip().lower()
            if val in ['1', 'true', 'yes', '是', 'y']:
                use_spac = True
            else:
                try:
                    use_spac = float(val) == 1.0
                except Exception:
                    use_spac = False
        else:
            try:
                use_spac = float(spac_flag) == 1.0
            except Exception:
                use_spac = False

        src = slide_spac if use_spac else slide_ipo
        new_slide = clone_slide(prs, src)
        mapping = row.copy()
        for key in [
            'company_intro',
            'company_name',
            'funding_2020',
            'funding_2021',
            'funding_2022',
            'funding_2023',
            'funding_2024',
            'ipo_date',
            'ipo_fund',
            'latest_price',
            'market_cap',
            'market_cap_level',
            'merger_valuation',
            'pe_ttm',
            'profit_2024',
            'ps_ttm',
            'revenue_2024',
            'spac_size',
            'ticker',
        ]:
            mapping.setdefault(key, '')
        replace_text_in_slide(new_slide, mapping)
        replace_logo(new_slide, row.get('logo_url', ''), cache)

    summary_slide = clone_slide(prs, slide_summary)
    replace_summary_image(summary_slide, chart_path)

    clone_slide(prs, slide_thanks)

    sldIdLst = prs.slides._sldIdLst
    for i in reversed(range(original_slide_count)):
        sldIdLst.remove(sldIdLst[i])

    prs.save(output_path)

    if cleanup_chart:
        try:
            os.remove(chart_path)
        except Exception:
            pass

    return output_path


def main():
    generate_ppt(
        excel_data_path=EXCEL,
        output_path=OUTPUT,
        template_path=TEMPLATE,
        chart_path=CHART_PATH,
    )


if __name__ == '__main__':
    main()
